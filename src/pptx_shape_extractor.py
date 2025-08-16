"""
PowerPoint Shape Extractor - Shape-aware parsing for presentations
Extracts and preserves relationships in organizational charts and diagrams
"""

import json
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
import structlog
import re

# Graceful import of python-pptx
try:
    from pptx import Presentation
    from pptx.shapes.group import GroupShape
    from pptx.shapes.connector import Connector
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    # Create mock classes for testing when python-pptx is not available
    PPTX_AVAILABLE = False
    
    class MockPresentation:
        def __init__(self, *args, **kwargs):
            self.slides = []
    
    class MockGroupShape:
        def __init__(self):
            self.shapes = []
            self.shape_id = 0
    
    class MockShape:
        def __init__(self):
            self.text_frame = None
            self.shape_type = 1
            self.shape_id = 0
            self.left = 0
            self.top = 0
            self.width = 0
            self.height = 0
    
    class MockMSOShapeType:
        AUTO_SHAPE = 1
        TEXT_BOX = 17
        GROUP = 6
        PICTURE = 13
        PLACEHOLDER = 14
        TABLE = 19
        CHART = 3
        DIAGRAM = 21
    
    Presentation = MockPresentation
    GroupShape = MockGroupShape
    Shape = MockShape
    MSO_SHAPE_TYPE = MockMSOShapeType()

logger = structlog.get_logger()

@dataclass
class ShapeElement:
    """Individual shape element from PowerPoint"""
    text: str
    shape_type: str
    shape_id: int
    left: int
    top: int
    width: int
    height: int
    slide_num: int
    is_grouped: bool = False
    group_id: Optional[int] = None
    
    @property
    def center_x(self) -> float:
        return self.left + self.width / 2
    
    @property
    def center_y(self) -> float:
        return self.top + self.height / 2

@dataclass
class ShapeGroup:
    """Group of related shapes forming a logical unit"""
    shapes: List[ShapeElement]
    group_id: str
    group_type: str  # 'explicit' for PowerPoint groups, 'proximity' for detected
    confidence: float = 0.0
    
    @property
    def text_lines(self) -> List[str]:
        """Get text lines from all shapes in group"""
        lines = []
        # Sort shapes by vertical position (top to bottom)
        sorted_shapes = sorted(self.shapes, key=lambda s: s.top)
        for shape in sorted_shapes:
            if shape.text.strip():
                lines.append(shape.text.strip())
        return lines
    
    @property
    def bbox(self) -> Tuple[int, int, int, int]:
        """Bounding box of entire group"""
        if not self.shapes:
            return (0, 0, 0, 0)
        
        left = min(s.left for s in self.shapes)
        top = min(s.top for s in self.shapes)
        right = max(s.left + s.width for s in self.shapes)
        bottom = max(s.top + s.height for s in self.shapes)
        
        return (left, top, right - left, bottom - top)

@dataclass
class PPTXOrganizationalUnit:
    """Organizational information extracted from PowerPoint shapes"""
    name: str
    title: str
    department: str
    shape_group_id: str
    confidence: float
    slide_number: int
    warnings: List[str]

class PPTXShapeExtractor:
    """
    PowerPoint shape extractor for organizational charts and structured diagrams
    Preserves shape relationships and extracts organizational information
    """
    
    def __init__(self):
        self.proximity_threshold = 100  # EMUs (English Metric Units) for shape proximity
        self.min_group_size = 2
        self.max_group_size = 10
    
    def extract_organizational_data(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract organizational data from PowerPoint preserving shape relationships
        """
        try:
            # Check if python-pptx is available
            if not PPTX_AVAILABLE:
                logger.warning("python-pptx not available, PowerPoint shape extraction disabled")
                return self._create_error_result("python-pptx package not installed")
            
            logger.info(f"Starting PowerPoint shape extraction: {pptx_path}")
            
            # Step 1: Extract all shapes with their properties
            shape_elements = self._extract_shape_elements(pptx_path)
            
            if not shape_elements:
                return self._create_empty_result("No shapes found in presentation")
            
            logger.info(f"Extracted {len(shape_elements)} shape elements")
            
            # Step 2: Identify and process shape groups
            shape_groups = self._identify_shape_groups(shape_elements)
            
            logger.info(f"Identified {len(shape_groups)} shape groups")
            
            # Step 3: Parse organizational information from groups
            org_units = self._parse_organizational_groups(shape_groups)
            
            logger.info(f"Extracted {len(org_units)} organizational units")
            
            # Step 4: Create structured result
            result = self._create_structured_result(org_units, shape_groups, shape_elements)
            
            logger.info("PowerPoint shape extraction completed successfully")
            return result
            
        except Exception as e:
            logger.error(f"PowerPoint shape extraction failed: {e}")
            return self._create_error_result(str(e))
    
    def _extract_shape_elements(self, pptx_path: str) -> List[ShapeElement]:
        """Extract all shapes from PowerPoint with their properties"""
        elements = []
        
        try:
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Process regular shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            elements.append(ShapeElement(
                                text=text,
                                shape_type=self._get_shape_type_name(shape),
                                shape_id=shape.shape_id,
                                left=shape.left,
                                top=shape.top,
                                width=shape.width,
                                height=shape.height,
                                slide_num=slide_num,
                                is_grouped=False
                            ))
                    
                    # Process grouped shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        group_elements = self._extract_from_group(shape, slide_num)
                        elements.extend(group_elements)
            
        except Exception as e:
            logger.error(f"Failed to extract shape elements: {e}")
        
        return elements
    
    def _extract_from_group(self, group_shape: GroupShape, slide_num: int) -> List[ShapeElement]:
        """Extract shapes from a PowerPoint group"""
        elements = []
        group_id = group_shape.shape_id
        
        try:
            for shape in group_shape.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        elements.append(ShapeElement(
                            text=text,
                            shape_type=self._get_shape_type_name(shape),
                            shape_id=shape.shape_id,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                            slide_num=slide_num,
                            is_grouped=True,
                            group_id=group_id
                        ))
                
                # Recursively process nested groups
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    nested_elements = self._extract_from_group(shape, slide_num)
                    elements.extend(nested_elements)
        
        except Exception as e:
            logger.warning(f"Failed to extract from group: {e}")
        
        return elements
    
    def _identify_shape_groups(self, elements: List[ShapeElement]) -> List[ShapeGroup]:
        """Identify shape groups - both explicit PowerPoint groups and proximity-based"""
        groups = []
        
        # Step 1: Process explicit PowerPoint groups
        grouped_elements = {}
        ungrouped_elements = []
        
        for element in elements:
            if element.is_grouped:
                group_key = f"pptx_group_{element.group_id}_slide_{element.slide_num}"
                if group_key not in grouped_elements:
                    grouped_elements[group_key] = []
                grouped_elements[group_key].append(element)
            else:
                ungrouped_elements.append(element)
        
        # Create ShapeGroup objects for explicit groups
        for group_key, shapes in grouped_elements.items():
            if len(shapes) >= self.min_group_size:
                groups.append(ShapeGroup(
                    shapes=shapes,
                    group_id=group_key,
                    group_type='explicit',
                    confidence=0.9  # High confidence for explicit groups
                ))
        
        # Step 2: Detect proximity-based groups for ungrouped shapes
        proximity_groups = self._detect_proximity_groups(ungrouped_elements)
        groups.extend(proximity_groups)
        
        return groups
    
    def _detect_proximity_groups(self, elements: List[ShapeElement]) -> List[ShapeGroup]:
        """Detect groups based on shape proximity"""
        groups = []
        used_indices = set()
        
        # Group shapes by slide first
        slide_elements = {}
        for i, element in enumerate(elements):
            if element.slide_num not in slide_elements:
                slide_elements[element.slide_num] = []
            slide_elements[element.slide_num].append((i, element))
        
        # Process each slide
        for slide_num, slide_shapes in slide_elements.items():
            for i, (idx1, shape1) in enumerate(slide_shapes):
                if idx1 in used_indices:
                    continue
                
                # Start new proximity group
                group_shapes = [shape1]
                used_indices.add(idx1)
                
                # Find nearby shapes
                for j, (idx2, shape2) in enumerate(slide_shapes):
                    if i == j or idx2 in used_indices:
                        continue
                    
                    # Check proximity to any shape in current group
                    for group_shape in group_shapes:
                        if self._shapes_are_proximate(group_shape, shape2):
                            group_shapes.append(shape2)
                            used_indices.add(idx2)
                            break
                
                # Create group if it has enough shapes
                if len(group_shapes) >= self.min_group_size and len(group_shapes) <= self.max_group_size:
                    group_id = f"proximity_group_{len(groups)}_slide_{slide_num}"
                    confidence = self._calculate_group_confidence(group_shapes)
                    
                    groups.append(ShapeGroup(
                        shapes=group_shapes,
                        group_id=group_id,
                        group_type='proximity',
                        confidence=confidence
                    ))
        
        return groups
    
    def _shapes_are_proximate(self, shape1: ShapeElement, shape2: ShapeElement) -> bool:
        """Check if two shapes are close enough to be grouped"""
        # Calculate distance between shape centers
        dx = shape1.center_x - shape2.center_x
        dy = shape1.center_y - shape2.center_y
        distance = (dx * dx + dy * dy) ** 0.5
        
        # Convert threshold to shape units (EMUs)
        threshold_emu = self.proximity_threshold * 914400  # 914400 EMUs per inch
        
        return distance <= threshold_emu
    
    def _calculate_group_confidence(self, shapes: List[ShapeElement]) -> float:
        """Calculate confidence score for a shape group"""
        if len(shapes) < 2:
            return 0.0
        
        confidence = 0.5  # Base confidence
        
        # Bonus for vertical or horizontal alignment
        x_positions = [s.center_x for s in shapes]
        y_positions = [s.center_y for s in shapes]
        
        x_variance = max(x_positions) - min(x_positions)
        y_variance = max(y_positions) - min(y_positions)
        
        # Check for good alignment
        if x_variance < 1000000:  # Well aligned horizontally (in EMUs)
            confidence += 0.2
        if y_variance < 1000000:  # Well aligned vertically
            confidence += 0.2
        
        # Bonus for typical org chart group size (2-4 shapes)
        if 2 <= len(shapes) <= 4:
            confidence += 0.1
        
        return min(confidence, 1.0)
    
    def _parse_organizational_groups(self, groups: List[ShapeGroup]) -> List[PPTXOrganizationalUnit]:
        """Parse shape groups to extract organizational information"""
        org_units = []
        
        for group in groups:
            if group.confidence < 0.5:  # Skip low confidence groups
                continue
            
            try:
                org_unit = self._parse_single_group(group)
                if org_unit:
                    org_units.append(org_unit)
            except Exception as e:
                logger.warning(f"Failed to parse group {group.group_id}: {e}")
        
        return org_units
    
    def _parse_single_group(self, group: ShapeGroup) -> Optional[PPTXOrganizationalUnit]:
        """Parse a single shape group to extract organizational info"""
        lines = group.text_lines
        
        if len(lines) < 2:
            return None
        
        name = ""
        title = ""
        department = ""
        warnings = []
        
        # Identify each line type
        for line in lines:
            line_type = self._identify_line_type(line)
            
            if line_type == "name" and not name:
                name = line
            elif line_type == "title" and not title:
                title = line
            elif line_type == "department" and not department:
                department = line
        
        # Fallback parsing patterns
        if not all([name, title]) and len(lines) >= 2:
            # Common patterns in org charts
            if len(lines) >= 3:
                # Pattern: Name, Title, Department
                if not name:
                    name = lines[0]
                if not title:
                    title = lines[1]
                if not department:
                    department = lines[2]
            elif len(lines) == 2:
                # Pattern: Name, Title
                if not name:
                    name = lines[0]
                if not title:
                    title = lines[1]
            
            warnings.append("Used fallback parsing pattern")
        
        # Validation
        if not name or not title:
            return None
        
        # Get slide number from first shape
        slide_num = group.shapes[0].slide_num if group.shapes else 0
        
        confidence = group.confidence
        if warnings:
            confidence *= 0.8
        
        return PPTXOrganizationalUnit(
            name=name.strip(),
            title=title.strip(),
            department=department.strip() if department else "Not specified",
            shape_group_id=group.group_id,
            confidence=confidence,
            slide_number=slide_num,
            warnings=warnings
        )
    
    def _identify_line_type(self, line: str) -> str:
        """Identify if a line contains name, title, or department"""
        line_lower = line.lower()
        
        # Department patterns
        dept_keywords = ['finance', 'accounting', 'hr', 'human resources', 'operations',
                        'marketing', 'sales', 'it', 'technology', 'legal', 'department',
                        'division', 'team', 'group']
        
        if any(keyword in line_lower for keyword in dept_keywords):
            return "department"
        
        # Title patterns
        title_keywords = ['cfo', 'ceo', 'cto', 'coo', 'director', 'manager', 'analyst',
                         'controller', 'treasurer', 'officer', 'senior', 'junior', 'lead',
                         'chief', 'president', 'vice', 'assistant', 'associate', 'coordinator',
                         'supervisor', 'head', 'specialist', 'executive']
        
        if any(keyword in line_lower for keyword in title_keywords):
            return "title"
        
        # Name patterns
        name_patterns = [
            r'^[A-Z][a-z]+ [A-Z][a-z]+$',  # First Last
            r'^[A-Z][a-z]+ [A-Z]\. [A-Z][a-z]+$',  # First M. Last
            r'^[A-Z][a-z]+ [A-Z][a-z]+ [A-Z][a-z]+$',  # First Middle Last
            r'^Dr\. [A-Z][a-z]+ [A-Z][a-z]+$',  # Dr. First Last
            r'^[A-Z][a-z]+, [A-Z][a-z]+$'  # Last, First
        ]
        
        for pattern in name_patterns:
            if re.match(pattern, line):
                return "name"
        
        # Default to title if unclear
        return "title"
    
    def _get_shape_type_name(self, shape) -> str:
        """Get human-readable shape type name"""
        try:
            shape_type_map = {
                MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
                MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
                MSO_SHAPE_TYPE.GROUP: "Group",
                MSO_SHAPE_TYPE.PICTURE: "Picture",
                MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
                MSO_SHAPE_TYPE.TABLE: "Table",
                MSO_SHAPE_TYPE.CHART: "Chart",
                MSO_SHAPE_TYPE.DIAGRAM: "Diagram"
            }
            
            return shape_type_map.get(shape.shape_type, "Unknown")
        except:
            return "Unknown"
    
    def _create_structured_result(self, org_units: List[PPTXOrganizationalUnit],
                                shape_groups: List[ShapeGroup],
                                elements: List[ShapeElement]) -> Dict[str, Any]:
        """Create final structured result"""
        
        return {
            "extraction_method": "pptx_shapes",
            "success": True,
            "organizational_units": [
                {
                    "name": unit.name,
                    "title": unit.title,
                    "department": unit.department,
                    "confidence": unit.confidence,
                    "slide_number": unit.slide_number,
                    "shape_group_id": unit.shape_group_id,
                    "warnings": unit.warnings
                }
                for unit in org_units
            ],
            "metadata": {
                "total_shapes": len(elements),
                "total_groups": len(shape_groups),
                "successful_extractions": len(org_units),
                "slides_processed": len(set(elem.slide_num for elem in elements)),
                "average_confidence": sum(unit.confidence for unit in org_units) / len(org_units) if org_units else 0.0
            },
            "shape_data": {
                "groups": [
                    {
                        "id": group.group_id,
                        "type": group.group_type,
                        "text_lines": group.text_lines,
                        "confidence": group.confidence,
                        "shape_count": len(group.shapes)
                    }
                    for group in shape_groups
                ]
            }
        }
    
    def _create_empty_result(self, reason: str) -> Dict[str, Any]:
        """Create empty result structure"""
        return {
            "extraction_method": "pptx_shapes",
            "success": False,
            "organizational_units": [],
            "metadata": {
                "total_shapes": 0,
                "total_groups": 0,
                "successful_extractions": 0,
                "slides_processed": 0,
                "average_confidence": 0.0,
                "error": reason
            },
            "shape_data": {"groups": []}
        }
    
    def _create_error_result(self, error: str) -> Dict[str, Any]:
        """Create error result structure"""
        return {
            "extraction_method": "pptx_shapes",
            "success": False,
            "organizational_units": [],
            "metadata": {
                "total_shapes": 0,
                "total_groups": 0,
                "successful_extractions": 0,
                "slides_processed": 0,
                "average_confidence": 0.0,
                "error": error
            },
            "shape_data": {"groups": []}
        }

# Convenience function
async def extract_pptx_organizational_data(pptx_path: str) -> Dict[str, Any]:
    """
    Async wrapper for PowerPoint organizational data extraction
    """
    if not PPTX_AVAILABLE:
        return {
            "extraction_method": "pptx_shapes",
            "success": False,
            "organizational_units": [],
            "metadata": {
                "total_shapes": 0,
                "total_groups": 0,
                "successful_extractions": 0,
                "slides_processed": 0,
                "average_confidence": 0.0,
                "error": "python-pptx package not installed"
            },
            "shape_data": {"groups": []}
        }
    
    extractor = PPTXShapeExtractor()
    return extractor.extract_organizational_data(pptx_path)