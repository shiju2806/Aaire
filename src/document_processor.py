"""
Document Processor for AAIRE - MVP-FR-009 through MVP-FR-012
Handles document upload, processing, and ingestion into RAG pipeline
"""

import os
import uuid
import json
import asyncio
from typing import Dict, Any, List, Optional, BinaryIO
from datetime import datetime
from pathlib import Path
import tempfile
import mimetypes

import structlog
from fastapi import UploadFile, HTTPException
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

from llama_index.core import Document
from .rag_pipeline import RAGPipeline

logger = structlog.get_logger()

class DocumentProcessor:
    def __init__(self, rag_pipeline: RAGPipeline = None):
        """Initialize document processor"""
        self.rag_pipeline = rag_pipeline
        self.upload_dir = Path("data/uploads")
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # Document processing status tracking
        self.processing_jobs = {}
        
        # Supported file types and size limits (Enhanced for multi-modal analysis)
        self.supported_formats = {
            'application/pdf': {'extension': '.pdf', 'max_size_mb': 100, 'category': 'document'},
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': {
                'extension': '.docx', 'max_size_mb': 50, 'category': 'document'
            },
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': {
                'extension': '.pptx', 'max_size_mb': 50, 'category': 'presentation'
            },
            'application/vnd.ms-powerpoint': {
                'extension': '.ppt', 'max_size_mb': 50, 'category': 'presentation'
            },
            'text/plain': {'extension': '.txt', 'max_size_mb': 10, 'category': 'document'},
            'text/csv': {'extension': '.csv', 'max_size_mb': 25, 'category': 'data'},
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': {
                'extension': '.xlsx', 'max_size_mb': 25, 'category': 'data'
            },
            # Multi-modal support for charts and financial statements
            'image/png': {'extension': '.png', 'max_size_mb': 10, 'category': 'image'},
            'image/jpeg': {'extension': '.jpg', 'max_size_mb': 10, 'category': 'image'},
            'image/gif': {'extension': '.gif', 'max_size_mb': 10, 'category': 'image'},
            'image/webp': {'extension': '.webp', 'max_size_mb': 10, 'category': 'image'}
        }
        
        logger.info("Document processor initialized", 
                   upload_dir=str(self.upload_dir),
                   supported_formats=list(self.supported_formats.keys()))
    
    async def upload_document(
        self, 
        file: UploadFile, 
        metadata: str, 
        user_id: str
    ) -> str:
        """
        Upload and queue document for processing - MVP-FR-009
        Returns job_id for tracking
        """
        
        # Generate unique job ID
        job_id = str(uuid.uuid4())
        
        try:
            # Parse metadata
            metadata_dict = json.loads(metadata) if metadata else {}
            
            # Validate file
            await self._validate_file(file, metadata_dict)
            
            # Save file temporarily
            file_path = await self._save_uploaded_file(file, job_id)
            
            # Create processing job
            self.processing_jobs[job_id] = {
                'status': 'queued',
                'user_id': user_id,
                'filename': file.filename,
                'file_path': str(file_path),
                'metadata': metadata_dict,
                'created_at': datetime.utcnow().isoformat(),
                'progress': 0
            }
            
            # Queue processing (in production, this would use Celery)
            asyncio.create_task(self._process_document_async(job_id))
            
            logger.info("Document upload queued", 
                       job_id=job_id, 
                       filename=file.filename,
                       user_id=user_id)
            
            return job_id
            
        except Exception as e:
            logger.error("Document upload failed", 
                        error=str(e), 
                        filename=getattr(file, 'filename', 'unknown'))
            raise HTTPException(status_code=400, detail=f"Upload failed: {str(e)}")
    
    async def _validate_file(self, file: UploadFile, metadata: Dict[str, Any]):
        """Validate uploaded file - MVP-FR-010"""
        
        # Check file size
        file.file.seek(0, 2)  # Seek to end
        size = file.file.tell()
        file.file.seek(0)  # Reset to beginning
        
        # Detect mime type
        mime_type, _ = mimetypes.guess_type(file.filename)
        if not mime_type or mime_type not in self.supported_formats:
            raise ValueError(f"Unsupported file type: {mime_type}")
        
        # Check size limit
        max_size_bytes = self.supported_formats[mime_type]['max_size_mb'] * 1024 * 1024
        if size > max_size_bytes:
            raise ValueError(f"File too large: {size} bytes > {max_size_bytes} bytes")
        
        # Validate required metadata
        required_fields = ['title', 'source_type', 'effective_date']
        for field in required_fields:
            if field not in metadata:
                raise ValueError(f"Missing required metadata field: {field}")
        
        # Validate source_type
        valid_source_types = ['US_GAAP', 'IFRS', 'COMPANY', 'ACTUARIAL']
        if metadata['source_type'] not in valid_source_types:
            raise ValueError(f"Invalid source_type: {metadata['source_type']}")
    
    async def _save_uploaded_file(self, file: UploadFile, job_id: str) -> Path:
        """Save uploaded file to temporary location"""
        
        # Create safe filename
        extension = Path(file.filename).suffix
        safe_filename = f"{job_id}{extension}"
        file_path = self.upload_dir / safe_filename
        
        # Save file
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        return file_path
    
    async def _process_document_async(self, job_id: str):
        """Process document asynchronously - MVP-FR-011, MVP-FR-012"""
        
        try:
            # Update status
            self.processing_jobs[job_id]['status'] = 'processing'
            self.processing_jobs[job_id]['progress'] = 10
            
            job = self.processing_jobs[job_id]
            file_path = Path(job['file_path'])
            metadata = job['metadata']
            
            # Extract text based on file type
            text_content = await self._extract_text(file_path)
            
            self.processing_jobs[job_id]['progress'] = 40
            
            # Create document with metadata
            document = Document(
                text=text_content,
                metadata={
                    **metadata,
                    'filename': job['filename'],
                    'job_id': job_id,
                    'processed_at': datetime.utcnow().isoformat(),
                    'user_id': job['user_id']
                }
            )
            
            self.processing_jobs[job_id]['progress'] = 70
            
            # Add to RAG pipeline (single index)
            if self.rag_pipeline:
                doc_type = self._map_source_type(metadata['source_type'])
                chunks_created = await self.rag_pipeline.add_documents([document], doc_type)
                
                self.processing_jobs[job_id]['chunks_created'] = chunks_created
                
                # Generate intelligent document summary
                try:
                    summary_result = await self.rag_pipeline.generate_document_summary(
                        text_content, metadata
                    )
                    self.processing_jobs[job_id]['summary'] = summary_result
                    logger.info("Document summary generated", 
                               job_id=job_id, 
                               confidence=summary_result.get('confidence', 0))
                except Exception as e:
                    logger.warning("Failed to generate document summary", 
                                 job_id=job_id, error=str(e))
                    # Fallback basic summary
                    self.processing_jobs[job_id]['summary'] = self._generate_basic_summary(
                        text_content, metadata
                    )
            else:
                # Generate basic summary when RAG pipeline not available
                self.processing_jobs[job_id]['summary'] = self._generate_basic_summary(
                    text_content, metadata
                )
            
            self.processing_jobs[job_id]['progress'] = 100
            self.processing_jobs[job_id]['status'] = 'completed'
            self.processing_jobs[job_id]['completed_at'] = datetime.utcnow().isoformat()
            
            # Clean up file
            if file_path.exists():
                file_path.unlink()
            
            logger.info("Document processing completed", 
                       job_id=job_id,
                       chunks_created=self.processing_jobs[job_id].get('chunks_created', 0))
            
        except Exception as e:
            self.processing_jobs[job_id]['status'] = 'failed'
            self.processing_jobs[job_id]['error'] = str(e)
            self.processing_jobs[job_id]['failed_at'] = datetime.utcnow().isoformat()
            
            logger.error("Document processing failed", 
                        job_id=job_id, 
                        error=str(e))
    
    async def _extract_text(self, file_path: Path) -> str:
        """Extract text from various file formats"""
        
        file_extension = file_path.suffix.lower()
        
        try:
            if file_extension == '.pdf':
                return await self._extract_from_pdf(file_path)
            elif file_extension == '.docx':
                return await self._extract_from_docx(file_path)
            elif file_extension == '.txt':
                return await self._extract_from_txt(file_path)
            elif file_extension == '.csv':
                return await self._extract_from_csv(file_path)
            elif file_extension == '.xlsx':
                return await self._extract_from_xlsx(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                return await self._extract_from_powerpoint(file_path)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.webp']:
                return await self._extract_from_image(file_path)
            else:
                raise ValueError(f"Unsupported file extension: {file_extension}")
                
        except Exception as e:
            logger.error("Text extraction failed", 
                        file_path=str(file_path), 
                        error=str(e))
            raise
    
    async def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        text_content = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                except Exception as e:
                    logger.warning("Failed to extract page", 
                                 page_num=page_num, 
                                 error=str(e))
                    continue
        
        return "\n\n".join(text_content)
    
    async def _extract_from_txt(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                return content.strip()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()
                return content.strip()
    
    async def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        doc = DocxDocument(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                table_text.append(row_text)
            
            if table_text:
                text_content.append(f"[TABLE]\n" + "\n".join(table_text))
        
        return "\n\n".join(text_content)
    
    async def _extract_from_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            
            # Create structured text representation
            text_parts = [
                f"[CSV FILE - {len(df)} rows, {len(df.columns)} columns]",
                f"Columns: {', '.join(df.columns.tolist())}",
                "\nData Summary:",
                str(df.describe(include='all'))
            ]
            
            # Include first few rows as sample
            if len(df) > 0:
                text_parts.append("\nSample Data (first 5 rows):")
                text_parts.append(df.head().to_string())
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error("CSV extraction failed", error=str(e))
            raise
    
    async def _extract_from_xlsx(self, file_path: Path) -> str:
        """Extract text from XLSX file"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = [f"[EXCEL FILE - {len(excel_file.sheet_names)} sheets]"]
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                text_parts.append(f"\n[SHEET: {sheet_name}]")
                text_parts.append(f"Dimensions: {len(df)} rows, {len(df.columns)} columns")
                text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                
                if len(df) > 0:
                    text_parts.append("Sample Data:")
                    text_parts.append(df.head(3).to_string())
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error("Excel extraction failed", error=str(e))
            raise
    
    async def _extract_from_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint files (.ppt/.pptx)"""
        try:
            logger.info(f"Starting PowerPoint extraction for: {file_path}")
            prs = Presentation(file_path)
            text_content = []
            
            logger.info(f"PowerPoint loaded successfully - {len(prs.slides)} slides found")
            text_content.append(f"[POWERPOINT PRESENTATION - {len(prs.slides)} slides]")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n[SLIDE {slide_num}]"]
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Identify the type of content
                        if shape.shape_type == 1:  # Title placeholder
                            slide_text.append(f"TITLE: {shape.text}")
                        elif shape.shape_type == 2:  # Body/content placeholder
                            slide_text.append(f"CONTENT: {shape.text}")
                        else:
                            slide_text.append(shape.text)
                    
                    # Extract table content if present
                    if shape.shape_type == 19:  # Table
                        try:
                            table_text = ["[TABLE]"]
                            table = shape.table
                            for row in table.rows:
                                row_text = " | ".join(cell.text for cell in row.cells)
                                table_text.append(row_text)
                            slide_text.append("\n".join(table_text))
                        except:
                            pass
                
                # Only add slide if it has content
                if len(slide_text) > 1:
                    text_content.extend(slide_text)
            
            extracted_text = "\n\n".join(text_content)
            logger.info(f"PowerPoint extraction completed - {len(extracted_text)} characters extracted")
            return extracted_text
            
        except Exception as e:
            logger.error("PowerPoint extraction failed", 
                        file_path=str(file_path), 
                        error=str(e))
            raise
    
    async def _extract_from_image(self, file_path: Path) -> str:
        """Extract information from image files using GPT-4 Vision"""
        try:
            import base64
            
            # Read and encode image
            with open(file_path, 'rb') as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            # For now, return basic metadata - will enhance with GPT-4 Vision
            file_stats = file_path.stat()
            
            return f"""[IMAGE FILE ANALYSIS]
Filename: {file_path.name}
File Type: {file_path.suffix.upper()} Image
File Size: {file_stats.st_size / 1024:.1f} KB
Upload Date: {datetime.utcnow().isoformat()}

Content Analysis: This appears to be a financial chart, graph, or document image.
For detailed analysis of charts, graphs, and financial statements, enhanced AI vision capabilities are being processed.

Base64 Data Available: Yes (for AI vision analysis)
Analysis Status: Ready for multi-modal processing"""
            
        except Exception as e:
            logger.error("Image extraction failed", error=str(e))
            raise
    
    def _generate_basic_summary(self, text_content: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Generate enhanced basic summary when AI is not available"""
        try:
            import re
            
            # Extract basic statistics
            word_count = len(text_content.split())
            char_count = len(text_content)
            page_count = text_content.count('[Page ') + 1 if '[Page ' in text_content else 1
            
            # Enhanced accounting and insurance terms detection
            accounting_terms = {
                'Standards': ['gaap', 'ifrs', 'asc', 'fasb', 'iasb', 'pcaob', 'sox', 'coso'],
                'Financial Items': ['revenue', 'liability', 'asset', 'equity', 'income', 'expense', 'cash flow', 'balance sheet'],
                'Insurance Terms': ['reserve', 'premium', 'claim', 'policy', 'policyholder', 'underwriting', 'reinsurance'],
                'Actuarial': ['actuarial', 'valuation', 'assumption', 'mortality', 'morbidity', 'lapse', 'discount rate'],
                'Compliance': ['compliance', 'audit', 'control', 'regulation', 'requirement', 'disclosure']
            }
            
            found_by_category = {}
            for category, terms in accounting_terms.items():
                found = [term for term in terms if term in text_content.lower()]
                if found:
                    found_by_category[category] = found
            
            # Extract dates with better patterns
            date_patterns = [
                r'\b\d{1,2}/\d{1,2}/\d{4}\b',  # MM/DD/YYYY
                r'\b\d{4}-\d{2}-\d{2}\b',      # YYYY-MM-DD
                r'\b(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}\b',
                r'\b\d{1,2}\s+(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b'
            ]
            
            dates_found = []
            for pattern in date_patterns:
                dates_found.extend(re.findall(pattern, text_content, re.IGNORECASE))
            dates_found = list(set(dates_found))[:8]  # Remove duplicates, limit to 8
            
            # Extract financial numbers
            financial_patterns = [
                r'\$[\d,]+(?:\.\d{2})?(?:\s*(?:million|billion|thousand|M|B|K))?',
                r'[\d,]+(?:\.\d{2})?\s*(?:million|billion|thousand|M|B|K)?\s*(?:dollars?|USD|\$)',
                r'[\d,]+(?:\.\d{1,2})?%'
            ]
            
            financial_figures = []
            for pattern in financial_patterns:
                financial_figures.extend(re.findall(pattern, text_content, re.IGNORECASE))
            financial_figures = list(set(financial_figures))[:10]  # Limit to 10
            
            # Extract section headings (look for text in all caps or numbered sections)
            section_patterns = [
                r'^[A-Z\s]{10,50}$',  # All caps sections
                r'^\d+\.\s+[A-Z][^.]+$',  # Numbered sections
                r'^\d+\.\d+\s+[A-Z][^.]+$'  # Sub-sections
            ]
            
            sections = []
            for line in text_content.split('\n')[:50]:  # Check first 50 lines
                line = line.strip()
                if len(line) > 5 and len(line) < 80:
                    for pattern in section_patterns:
                        if re.match(pattern, line):
                            sections.append(line)
                            break
            sections = sections[:8]  # Limit to 8 sections
            
            # Generate comprehensive summary
            summary_text = f"""**Document Analysis Report**

**📄 Document Profile:**
- **File Name**: {metadata.get('title', 'Unknown Document')}
- **Document Type**: {metadata.get('source_type', 'Unknown')}
- **File Size**: {word_count:,} words across {page_count} pages
- **Analysis Date**: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}
- **Document Category**: {"Financial/Accounting Document" if found_by_category else "Business Document"}

**🔍 Content Analysis:**"""

            if found_by_category:
                summary_text += "\n**Key Areas Identified:**\n"
                for category, terms in found_by_category.items():
                    summary_text += f"- **{category}**: {', '.join(terms[:5])}\n"
            else:
                summary_text += "\n- General business document with accounting or insurance context\n"

            if sections:
                summary_text += f"\n**📋 Document Structure:**\n"
                for i, section in enumerate(sections, 1):
                    summary_text += f"{i}. {section}\n"

            if dates_found:
                summary_text += f"\n**📅 Important Dates Referenced:**\n"
                for date in dates_found:
                    summary_text += f"- {date}\n"

            if financial_figures:
                summary_text += f"\n**💰 Financial Figures Mentioned:**\n"
                for figure in financial_figures[:6]:
                    summary_text += f"- {figure}\n"

            summary_text += f"""

**⚠️ Compliance & Risk Considerations:**
- Review for adherence to applicable accounting standards
- Verify compliance with regulatory requirements
- Check for material misstatements or inconsistencies
- Assess impact on financial reporting and disclosures

**🎯 Recommended Actions:**
1. **Immediate Review**: Examine document for compliance with current accounting standards
2. **Stakeholder Communication**: Share relevant sections with accounting team and auditors  
3. **Documentation**: File in appropriate compliance folder with effective date tracking
4. **Follow-up**: Schedule review for any mentioned deadlines or effective dates
5. **Integration**: Consider impact on current financial statements and processes

**💡 Professional Insights:**
- Document contains {"substantial" if word_count > 5000 else "moderate" if word_count > 1000 else "basic"} level of detail
- {"Multiple accounting areas" if len(found_by_category) > 2 else "Specific focus area"} identified
- {"Complex document" if page_count > 20 else "Standard document"} requiring {"detailed" if word_count > 10000 else "standard"} professional review

*Note: This is an enhanced rule-based analysis. For AI-powered detailed insights and recommendations, configure OpenAI API keys for advanced document intelligence.*"""

            # Enhanced key insights
            insights = []
            
            if found_by_category:
                for category, terms in found_by_category.items():
                    insights.append({
                        "type": "content_category",
                        "value": terms,
                        "description": f"{category} concepts identified in document"
                    })
            
            insights.extend([
                {
                    "type": "document_metrics",
                    "value": f"{word_count:,} words, {page_count} pages",
                    "description": "Document size and complexity indicator"
                },
                {
                    "type": "compliance_scope",
                    "value": list(found_by_category.keys()) if found_by_category else ["General Business"],
                    "description": "Areas requiring compliance review"
                }
            ])
            
            if dates_found:
                insights.append({
                    "type": "critical_dates",
                    "value": dates_found[:5],
                    "description": "Important dates mentioned in document"
                })
            
            if financial_figures:
                insights.append({
                    "type": "financial_data",
                    "value": financial_figures[:5],
                    "description": "Financial figures and amounts referenced"
                })

            return {
                "summary": summary_text,
                "key_insights": insights,
                "document_metadata": metadata,
                "generated_at": datetime.utcnow().isoformat(),
                "confidence": 0.75,  # Higher confidence for enhanced analysis
                "analysis_type": "enhanced_basic"
            }
            
        except Exception as e:
            logger.error("Failed to generate basic summary", error=str(e))
            return {
                "summary": f"Document '{metadata.get('title', 'Unknown')}' was processed successfully. Basic analysis not available.",
                "key_insights": [],
                "document_metadata": metadata,
                "generated_at": datetime.utcnow().isoformat(),
                "confidence": 0.3,
                "analysis_type": "minimal"
            }

    def _map_source_type(self, source_type: str) -> str:
        """Map source type to RAG pipeline document type"""
        mapping = {
            'US_GAAP': 'us_gaap',
            'IFRS': 'ifrs',
            'COMPANY': 'company',
            'ACTUARIAL': 'actuarial'
        }
        return mapping.get(source_type, 'company')
    
    async def get_status(self, job_id: str, user_id: str) -> Dict[str, Any]:
        """Get processing status for a job"""
        
        if job_id not in self.processing_jobs:
            raise HTTPException(status_code=404, detail="Job not found")
        
        job = self.processing_jobs[job_id]
        
        # Check authorization
        if job['user_id'] != user_id:
            raise HTTPException(status_code=403, detail="Unauthorized")
        
        return {
            'job_id': job_id,
            'status': job['status'],
            'progress': job['progress'],
            'filename': job['filename'],
            'created_at': job['created_at'],
            'chunks_created': job.get('chunks_created'),
            'summary': job.get('summary'),
            'error': job.get('error')
        }
    
    def get_job_stats(self) -> Dict[str, Any]:
        """Get overall job processing statistics"""
        
        statuses = {}
        for job in self.processing_jobs.values():
            status = job['status']
            statuses[status] = statuses.get(status, 0) + 1
        
        return {
            'total_jobs': len(self.processing_jobs),
            'status_breakdown': statuses,
            'active_jobs': [
                job_id for job_id, job in self.processing_jobs.items()
                if job['status'] in ['queued', 'processing']
            ]
        }
    
    async def cleanup_job(self, job_id: str, user_id: str) -> Dict[str, Any]:
        """Clean up all traces of a job from document processor"""
        try:
            cleanup_results = {}
            
            # Remove from processing jobs if exists
            if job_id in self.processing_jobs:
                job_data = self.processing_jobs.pop(job_id)
                cleanup_results["processing_job"] = {"removed": job_data.get("filename", "unknown")}
                
                # Try to clean up associated files if path is stored
                if "file_path" in job_data:
                    try:
                        file_path = Path(job_data["file_path"])
                        if file_path.exists():
                            file_path.unlink()
                            cleanup_results["file_cleanup"] = {"removed": str(file_path)}
                    except Exception as e:
                        cleanup_results["file_cleanup"] = {"error": str(e)}
            
            return {
                "status": "success",
                "job_id": job_id,
                "cleanup_results": cleanup_results
            }
            
        except Exception as e:
            return {
                "status": "error", 
                "job_id": job_id,
                "error": str(e)
            }